#!/usr/bin/env python3
"""
Script para gerar apresentação em PowerPoint sobre o Sistema de Roteamento Multimodal do Porto
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Cores padrão
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 102, 204)
DARK_GREEN = RGBColor(0, 102, 51)
ORANGE = RGBColor(255, 102, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(80, 80, 80)

def add_title_slide(prs, title, subtitle, author=""):
    """Adiciona slide de título"""
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Fundo azul escuro
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitle
        subtitle_p.font.size = Pt(28)
        subtitle_p.font.color.rgb = LIGHT_BLUE
        subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Autor
    if author:
        author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
        author_frame = author_box.text_frame
        author_frame.word_wrap = True
        author_p = author_frame.paragraphs[0]
        author_p.text = author
        author_p.font.size = Pt(16)
        author_p.font.color.rgb = WHITE
        author_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    """Adiciona slide com conteúdo em lista"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Fundo branco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Barra azul no topo
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.color.rgb = DARK_BLUE
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Conteúdo
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.3))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        if isinstance(item, tuple):
            text, level = item
        else:
            text = item
            level = 0
        
        p.text = text
        p.level = level
        p.font.size = Pt(18) if level == 0 else Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Adiciona slide com duas colunas"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Fundo branco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Barra azul no topo
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = DARK_BLUE
    top_bar.line.color.rgb = DARK_BLUE
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Coluna esquerda - Título
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.2), Inches(0.4))
    left_title_frame = left_title_box.text_frame
    left_title_p = left_title_frame.paragraphs[0]
    left_title_p.text = left_title
    left_title_p.font.size = Pt(20)
    left_title_p.font.bold = True
    left_title_p.font.color.rgb = DARK_BLUE
    
    # Coluna esquerda - Conteúdo
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.2), Inches(4.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_items):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Coluna direita - Título
    right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.1), Inches(4.2), Inches(0.4))
    right_title_frame = right_title_box.text_frame
    right_title_p = right_title_frame.paragraphs[0]
    right_title_p.text = right_title
    right_title_p.font.size = Pt(20)
    right_title_p.font.bold = True
    right_title_p.font.color.rgb = DARK_BLUE
    
    # Coluna direita - Conteúdo
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.6), Inches(4.2), Inches(4.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_items):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    return slide

def create_presentation():
    """Cria a apresentação completa"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Título
    add_title_slide(
        prs,
        "Sistema de Roteamento Multimodal",
        "Otimização Multi-Objetivo para a Área Metropolitana do Porto",
        "Grupo 6 - CIN 2025\nUniversidade do Minho"
    )
    
    # Slide 2: Índice
    add_content_slide(prs, "Índice", [
        "Visão Geral do Projeto",
        "Arquitetura e Tecnologia",
        "Algoritmos Implementados",
        "Dados e Implementação",
        "Avaliação e Testes",
        "Resultados e Conclusões"
    ])
    
    # Slide 3: Visão Geral
    add_content_slide(prs, "Visão Geral do Projeto", [
        "🎯 Objetivo: Motor de roteamento multimodal inteligente",
        "Critérios de Otimização:",
        ("⏱  Tempo de viagem (minimizar)", 1),
        ("♻  Emissões de CO₂ (minimizar)", 1),
        ("🚶  Exercício físico (maximizar)", 1),
        "📊 Resultado: Fronteira de Pareto com múltiplas rotas eficientes"
    ])
    
    # Slide 4: Características Principais
    add_content_slide(prs, "Características Principais", [
        "✓ Otimização Multi-Objetivo com fronteira Pareto rigorosa",
        "✓ Dados Reais: GTFS (Metro + STCP) + OpenStreetMap",
        "✓ 3 Algoritmos Avançados:",
        ("A* Heurístico (rápido, 2-5s)", 1),
        ("Dijkstra Multi-Label (exaustivo, 100% garantido)", 1),
        ("ACO Estocástico (criativo, 10-15s)", 1),
        "✓ Análise Geográfica com ruas reais",
        "✓ 22 Casos de Teste (trivial a extremo)"
    ])
    
    # Slide 5: Stack Tecnológico
    add_content_slide(prs, "Stack Tecnológico", [
        "Python 3.12+ - Linguagem principal",
        "NetworkX 3.x - Estrutura de grafos",
        "OSMnx 1.x - Dados geográficos (OpenStreetMap)",
        "FastAPI 0.100+ - API REST",
        "Pandas 2.x - Processamento GTFS",
        "Folium 0.14+ - Visualização de mapas",
        "NumPy 1.24+ - Cálculos numéricos",
        "Poetry 1.5+ - Gestão de dependências"
    ])
    
    # Slide 6: Arquitetura
    add_two_column_slide(
        prs,
        "Arquitetura do Projeto",
        "Componentes",
        [
            "services/ - Lógica de negócio",
            "utils/ - Operações auxiliares",
            "models/ - Estruturas de dados",
            "algoritms/ - Implementações",
            "graph.py - Rede multimodal",
            "solution.py - Classe de solução"
        ],
        "Algoritmos",
        [
            "a_star.py - Heurístico rápido",
            "dijkstra.py - Exaustivo garantido",
            "aco.py - Estocástico criativo",
            "Cada retorna Fronteira Pareto",
            "Comparação automática",
            "Validação cruzada"
        ]
    )
    
    # Slide 7: A* Multi-Objetivo
    add_content_slide(prs, "A* Multi-Objetivo", [
        "Tipo: Algoritmo heurístico guloso",
        "Complexidade: O(E log V)",
        "Tempo: 2-5 segundos por consulta",
        "Características:",
        ("✓ Utiliza heurísticas para guiar a busca", 1),
        ("✓ Prioriza nós promissores", 1),
        ("✓ Encontra rotas de boa qualidade rapidamente", 1),
        ("✗ Não garante Fronteira Pareto completa", 1),
        "Ideal para: Aplicações em tempo real"
    ])
    
    # Slide 8: Dijkstra Multi-Label
    add_content_slide(prs, "Dijkstra Multi-Label (Exaustivo)", [
        "Tipo: Algoritmo de rótulo-setting multi-objetivo",
        "Complexidade: O(E × L log V) onde L = nº de etiquetas",
        "Tempo: 30-60 segundos por consulta",
        "Características:",
        ("✓ Mantém múltiplas etiquetas por nó", 1),
        ("✓ Propaga todas as soluções não-dominadas", 1),
        ("✓ Remove soluções dominadas iterativamente", 1),
        ("✓ GARANTE Fronteira de Pareto completa e exata", 1),
        "Ideal para: Gold standard de otimização multi-objetivo"
    ])
    
    # Slide 9: ACO
    add_content_slide(prs, "ACO (Ant Colony Optimization)", [
        "Tipo: Algoritmo estocástico inspirado na natureza",
        "Complexidade: O(I × A × E) - I iterações, A formigas",
        "Tempo: 10-15 segundos por consulta",
        "Características:",
        ("✓ Simula comportamento coletivo de formigas", 1),
        ("✓ Cada formiga constrói uma solução", 1),
        ("✓ Atualiza feromónios baseado em qualidade", 1),
        ("✗ Heurístico (pode não encontrar todas as soluções)", 1),
        "Ideal para: Exploração criativa e inovação"
    ])
    
    # Slide 10: Dados
    add_two_column_slide(
        prs,
        "Fontes de Dados",
        "GTFS",
        [
            "🚇 Metro do Porto",
            "• 95+ paragens",
            "• 6 linhas",
            "• Horários atualizados",
            "🚌 STCP (Autocarro)",
            "• 600+ paragens",
            "• Múltiplas linhas"
        ],
        "OpenStreetMap (OSM)",
        [
            "🗺 Rede viária completa",
            "🚶 Caminhos pedonais",
            "⏱ Velocidades estimadas",
            "📍 Coordenadas geográficas",
            "Cálculos Derivados:",
            "• CO₂, transferências, caminhada"
        ]
    )
    
    # Slide 11: Exemplo de Uso
    add_content_slide(prs, "Exemplo de Uso - Python API", [
        "from app.services.graph import GraphRoute",
        "from app.services.algoritms.a_star import optimized_multi_objective_routing",
        "",
        "graph = GraphRoute('Casa da Musica', 'Casino da Póvoa')",
        "routes = optimized_multi_objective_routing(...)",
        "",
        "for rota in routes:",
        "    print(f'Tempo: {rota.total_time} min')",
        "    print(f'CO2: {rota.total_co2}g')",
        "    print(f'Caminhada: {rota.total_walk_km}km')"
    ])
    
    # Slide 12: Avaliação
    add_content_slide(prs, "Framework de Avaliação", [
        "📊 Metodologia:",
        ("22 Casos de Teste (trivial a extremo)", 1),
        ("3 Algoritmos com comparação automática", 1),
        ("Métricas: Tempo, CO₂, Exercício físico", 1),
        ("Validação via pytest", 1),
        "",
        "Categorias:",
        ("Trivial: <1km | Moderado: 1-5km", 1),
        ("Desafiante: 5-20km | Extremo: múltiplas transferências", 1),
        "Execução: python -m app.test_cases"
    ])
    
    # Slide 13: Comparação de Algoritmos
    add_content_slide(prs, "Comparação de Algoritmos", [
        "⚡ Tempo de Execução: A* > ACO > Dijkstra",
        "✓ Qualidade Garantida: Dijkstra > ACO > A*",
        "📊 Cobertura Pareto: Dijkstra > ACO ≈ A*",
        "💡 Criatividade: ACO > A* > Dijkstra",
        "📈 Escalabilidade: A* > ACO > Dijkstra",
        "",
        "Recomendações:",
        ("Use Dijkstra para garantias de otimalidade", 1),
        ("Use A* para tempo real", 1),
        ("Use ACO para exploração criativa", 1)
    ])
    
    # Slide 14: Resultados Esperados
    add_content_slide(prs, "Exemplo de Resultados", [
        "Rota: Casa da Música → Casino da Póvoa de Varzim",
        "",
        "Rota 1 (Rápida):    45 min | 850g CO₂ | 2.5km caminhada",
        "Rota 2 (Ecológica): 58 min | 350g CO₂ | 5.2km caminhada",
        "Rota 3 (Activa):    62 min | 400g CO₂ | 8.1km caminhada",
        "Rota 4 (Balanceada):52 min | 550g CO₂ | 4.8km caminhada",
        "",
        "⚡ Nenhuma rota é superior em todos os critérios",
        "✓ Decisão depende de prioridades do utilizador"
    ])
    
    # Slide 15: Realizações
    add_two_column_slide(
        prs,
        "Principais Realizações",
        "Implementação",
        [
            "✓ 3 algoritmos funcionais",
            "✓ Integração GTFS completa",
            "✓ Grafo multimodal",
            "✓ 22 casos de teste",
            "✓ Framework avaliação"
        ],
        "Inovação",
        [
            "✓ Fronteira Pareto real",
            "✓ Multi-critério genuíno",
            "✓ Dados reais (GTFS)",
            "✓ 3 perspectivas diferentes",
            "✓ Interface amigável"
        ]
    )
    
    # Slide 16: Impacto
    add_content_slide(prs, "Impacto Potencial", [
        "🚗 Reduzir emissões de carbono",
        "   Oferecendo rotas ecológicas alternativas",
        "",
        "⏱ Otimizar tempo de deslocação",
        "   Através de roteamento inteligente",
        "",
        "🏃 Combater sedentarismo",
        "   Com opções que maximizam exercício físico",
        "",
        "Sustentabilidade urbana e bem-estar pessoal"
    ])
    
    # Slide 17: Desafios
    add_content_slide(prs, "Desafios Encontrados", [
        "🔧 Técnicos:",
        ("Complexidade de multi-objetivo | Integração GTFS", 1),
        "",
        "📋 Metodológicos:",
        ("Métricas CO₂ consistentes | Tempos realistas", 1),
        "",
        "✓ Soluções Implementadas:",
        ("Dijkstra multi-label robusto", 1),
        ("Validação cruzada entre algoritmos", 1),
        ("Testes extensivos (22 casos)", 1)
    ])
    
    # Slide 18: Trabalho Futuro
    add_content_slide(prs, "Trabalho Futuro", [
        "🌍 Expansão Geográfica",
        ("Outras cidades + bike-sharing", 1),
        "",
        "🔬 Melhorias Algorítmicas",
        ("NSGA-II | Otimização paralela", 1),
        "",
        "🎨 Interface Utilizador",
        ("Aplicação web interativa | Mapas 3D", 1),
        "",
        "♻ Sustentabilidade",
        ("Impacto real | Dados de energia", 1)
    ])
    
    # Slide 19: Conclusões
    add_content_slide(prs, "Conclusões", [
        "✅ Sistema robusto de roteamento multimodal implementado",
        "",
        "✅ Algoritmos multi-objetivo de vanguarda",
        "",
        "✅ Integração com dados reais (GTFS + OSM)",
        "",
        "✅ Fronteira Pareto rigorosa e validada",
        "",
        "✅ Demonstração de Computação Inspirada na Natureza",
        "   para mobilidade urbana sustentável"
    ])
    
    # Slide 20: Final
    add_title_slide(
        prs,
        "Obrigado!",
        "Questões?",
        "Repositório: github.com/MIA-CDFR/CIN_GRUPO6\nGrupo 6 - CIN 2025 - Universidade do Minho"
    )
    
    return prs

if __name__ == "__main__":
    print("🎬 Gerando apresentação PowerPoint...")
    prs = create_presentation()
    output_path = "relatorio_apresentacao.pptx"
    prs.save(output_path)
    print(f"✅ Apresentação salva em: {output_path}")
    print(f"📊 Total de slides: {len(prs.slides)}")
